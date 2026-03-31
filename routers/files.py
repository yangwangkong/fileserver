import hashlib
import mimetypes
import shutil
import subprocess
import tempfile
from html import escape
from pathlib import Path
from typing import Optional

from docx import Document as DocxDocument
from docx.oxml.ns import qn
from docx.table import Table
from docx.text.paragraph import Paragraph
from fastapi import APIRouter, File, HTTPException, Request, UploadFile
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse, RedirectResponse
from openpyxl import load_workbook
from pptx import Presentation

from config import APP_DIR, settings
from routers.auth import get_csrf_token, is_authenticated, validate_csrf
from tpl import templates

router = APIRouter()

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".svg", ".ico", ".tiff"}
VIDEO_EXTS = {".mp4", ".mkv", ".avi", ".mov", ".webm", ".m4v", ".flv", ".wmv", ".rmvb"}
AUDIO_EXTS = {".mp3", ".flac", ".wav", ".ogg", ".m4a", ".aac", ".opus", ".wma"}
TEXT_EXTS = {
    ".txt", ".md", ".py", ".js", ".ts", ".jsx", ".tsx", ".html", ".htm",
    ".css", ".json", ".yaml", ".yml", ".xml", ".csv", ".log", ".sh",
    ".bat", ".ini", ".cfg", ".conf", ".toml", ".env", ".gitignore",
    ".c", ".cpp", ".h", ".hpp", ".java", ".go", ".rs", ".rb", ".php",
    ".swift", ".kt", ".sql", ".r", ".lua", ".vim", ".dockerfile",
}
ARCHIVE_EXTS = {".zip", ".rar", ".7z", ".tar", ".gz", ".bz2", ".xz", ".tar.gz", ".tar.bz2"}
WORD_EXTS = {".doc", ".docx"}
EXCEL_EXTS = {".xls", ".xlsx", ".xlsm", ".xltx", ".xltm"}
POWERPOINT_EXTS = {".ppt", ".pptx", ".pptm"}
OFFICE_EXTS = WORD_EXTS | EXCEL_EXTS | POWERPOINT_EXTS

TEXT_PREVIEW_LIMIT = 1 * 1024 * 1024
OFFICE_SHEET_ROW_LIMIT = 60
OFFICE_SHEET_COL_LIMIT = 16
OFFICE_CACHE_DIR = Path(tempfile.gettempdir()) / "fileserver_office_preview"


def classify(ext: str) -> str:
    ext = ext.lower()
    if ext in IMAGE_EXTS:
        return "image"
    if ext in VIDEO_EXTS:
        return "video"
    if ext in AUDIO_EXTS:
        return "audio"
    if ext == ".pdf":
        return "pdf"
    if ext in WORD_EXTS:
        return "word"
    if ext in EXCEL_EXTS:
        return "excel"
    if ext in POWERPOINT_EXTS:
        return "powerpoint"
    if ext in TEXT_EXTS:
        return "text"
    if ext in ARCHIVE_EXTS:
        return "archive"
    return "other"


def get_mount(name: str):
    for mount in settings.mounts:
        if mount.name == name:
            return mount
    return None


def safe_resolve(base: str, rel: str = "") -> Optional[Path]:
    base_path = Path(base).resolve()
    try:
        target = (base_path / rel).resolve() if rel else base_path
        target.relative_to(base_path)
        return target
    except (ValueError, OSError):
        return None


def format_size(size: int) -> str:
    for unit in ("B", "KB", "MB", "GB", "TB"):
        if size < 1024:
            return f"{size:.1f} {unit}" if unit != "B" else f"{size} B"
        size /= 1024
    return f"{size:.1f} PB"


def office_cache_key(target: Path) -> str:
    target = target.resolve()
    stat = target.stat()
    payload = f"{target}|{stat.st_mtime_ns}|{stat.st_size}".encode("utf-8")
    return hashlib.sha1(payload).hexdigest()[:20]


def office_cache_dir(target: Path) -> Path:
    cache_dir = OFFICE_CACHE_DIR / office_cache_key(target)
    cache_dir.mkdir(parents=True, exist_ok=True)
    return cache_dir


def cleanup_cache_dir(cache_dir: Path) -> None:
    if cache_dir.exists():
        shutil.rmtree(cache_dir, ignore_errors=True)
    cache_dir.mkdir(parents=True, exist_ok=True)


def get_soffice_path() -> Optional[str]:
    candidates = [
        APP_DIR / "LibreOffice" / "program" / "soffice.exe",
        APP_DIR / "LibreOfficePortable" / "App" / "libreoffice" / "program" / "soffice.exe",
        shutil.which("soffice"),
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return str(candidate)
    return None


def convert_with_soffice(target: Path, out_dir: Path, convert_to: str, out_ext: str) -> Optional[Path]:
    soffice = get_soffice_path()
    if not soffice:
        return None

    subprocess.run(
        [soffice, "--headless", "--convert-to", convert_to, "--outdir", str(out_dir), str(target)],
        check=True,
        capture_output=True,
    )
    converted = out_dir / f"{target.stem}{out_ext}"
    return converted if converted.exists() else None


def convert_to_pdf_with_soffice(target: Path) -> Optional[Path]:
    cache_dir = office_cache_dir(target)
    pdf_path = cache_dir / f"{target.stem}.pdf"
    if pdf_path.exists():
        return pdf_path
    return convert_with_soffice(target, cache_dir, "pdf", ".pdf")


def build_office_pdf_preview(target: Path) -> Optional[str]:
    pdf_path = convert_to_pdf_with_soffice(target)
    if not pdf_path:
        return None
    return f"/office-cache/{office_cache_key(target)}/{pdf_path.name}"


def convert_with_word(target: Path, out_path: Path) -> Optional[Path]:
    try:
        import pythoncom
        from win32com.client import DispatchEx
    except Exception:  # noqa: BLE001
        return None

    pythoncom.CoInitialize()
    app = None
    document = None
    try:
        app = DispatchEx("Word.Application")
        app.Visible = False
        app.DisplayAlerts = 0
        document = app.Documents.Open(str(target), False, True, False)
        document.SaveAs2(str(out_path), FileFormat=16)
        return out_path if out_path.exists() else None
    finally:
        if document is not None:
            document.Close(False)
        if app is not None:
            app.Quit()
        pythoncom.CoUninitialize()


def convert_with_excel(target: Path, out_path: Path) -> Optional[Path]:
    try:
        import pythoncom
        from win32com.client import DispatchEx
    except Exception:  # noqa: BLE001
        return None

    pythoncom.CoInitialize()
    app = None
    workbook = None
    try:
        app = DispatchEx("Excel.Application")
        app.Visible = False
        app.DisplayAlerts = False
        workbook = app.Workbooks.Open(str(target), False, True)
        workbook.SaveAs(str(out_path), FileFormat=51)
        return out_path if out_path.exists() else None
    finally:
        if workbook is not None:
            workbook.Close(False)
        if app is not None:
            app.Quit()
        pythoncom.CoUninitialize()


def export_ppt_slides_with_com(target: Path, out_dir: Path) -> list[Path]:
    try:
        import pythoncom
        from win32com.client import DispatchEx
    except Exception:  # noqa: BLE001
        return []

    pythoncom.CoInitialize()
    app = None
    presentation = None
    try:
        app = DispatchEx("PowerPoint.Application")
        presentation = app.Presentations.Open(str(target), False, False, False)
        images: list[Path] = []
        for index, slide in enumerate(presentation.Slides, start=1):
            image_path = out_dir / f"slide-{index:03d}.png"
            slide.Export(str(image_path), "PNG", 1600, 900)
            if image_path.exists():
                images.append(image_path)
        return images
    finally:
        if presentation is not None:
            presentation.Close()
        if app is not None:
            app.Quit()
        pythoncom.CoUninitialize()


def convert_legacy_office(target: Path, kind: str) -> tuple[Path, Optional[str]]:
    ext = target.suffix.lower()
    if kind == "word" and ext == ".doc":
        cache_dir = office_cache_dir(target)
        out_path = cache_dir / f"{target.stem}.docx"
        if out_path.exists():
            return out_path, "已将 DOC 转换为 DOCX 后预览。"
        converted = convert_with_word(target, out_path) or convert_with_soffice(target, cache_dir, "docx", ".docx")
        if converted:
            return converted, "已将 DOC 转换为 DOCX 后预览。"
        raise RuntimeError("当前环境缺少 Microsoft Word 或 LibreOffice，无法预览 DOC 文件。")

    if kind == "excel" and ext == ".xls":
        cache_dir = office_cache_dir(target)
        out_path = cache_dir / f"{target.stem}.xlsx"
        if out_path.exists():
            return out_path, "已将 XLS 转换为 XLSX 后预览。"
        converted = convert_with_excel(target, out_path) or convert_with_soffice(target, cache_dir, "xlsx", ".xlsx")
        if converted:
            return converted, "已将 XLS 转换为 XLSX 后预览。"
        raise RuntimeError("当前环境缺少 Microsoft Excel 或 LibreOffice，无法预览 XLS 文件。")

    if kind == "powerpoint" and ext == ".ppt":
        cache_dir = office_cache_dir(target)
        out_path = cache_dir / f"{target.stem}.pptx"
        if out_path.exists():
            return out_path, "已将 PPT 转换为 PPTX 后预览。"
        converted = convert_with_soffice(target, cache_dir, "pptx", ".pptx")
        if converted:
            return converted, "已将 PPT 转换为 PPTX 后预览。"
        return target, "旧版 PPT 将尝试直接导出幻灯片图片。"

    return target, None


def render_docx_preview(target: Path) -> tuple[str, str]:
    document = DocxDocument(str(target))
    blocks: list[str] = []

    def render_runs(paragraph: Paragraph) -> str:
        parts: list[str] = []
        for run in paragraph.runs:
            text = run.text
            if not text:
                continue
            content = escape(text).replace("\n", "<br>")
            if run.bold:
                content = f"<strong>{content}</strong>"
            if run.italic:
                content = f"<em>{content}</em>"
            if run.underline:
                content = f"<span class='underline'>{content}</span>"
            parts.append(content)
        return "".join(parts) or escape(paragraph.text).replace("\n", "<br>")

    def render_paragraph(paragraph: Paragraph) -> None:
        text = paragraph.text.strip()
        if not text:
            return

        style_name = (paragraph.style.name or "").lower()
        content = render_runs(paragraph)
        if style_name in {"title", "subtitle"} or "heading 1" in style_name:
            blocks.append(f"<h1 class='text-2xl font-bold text-gray-900 mt-2 mb-4'>{content}</h1>")
        elif "heading 2" in style_name:
            blocks.append(f"<h2 class='text-xl font-semibold text-gray-900 mt-5 mb-3'>{content}</h2>")
        elif "heading 3" in style_name:
            blocks.append(f"<h3 class='text-lg font-semibold text-gray-800 mt-4 mb-2'>{content}</h3>")
        elif "list" in style_name:
            blocks.append(f"<p class='my-2 pl-5 text-gray-700 leading-7'>• {content}</p>")
        else:
            blocks.append(f"<p class='my-2 text-gray-700 leading-7'>{content}</p>")

    def render_table(table: Table) -> None:
        rows_html: list[str] = []
        for row_index, row in enumerate(table.rows):
            cell_tag = "th" if row_index == 0 else "td"
            cells_html: list[str] = []
            for cell in row.cells:
                pieces = [paragraph.text.strip() for paragraph in cell.paragraphs if paragraph.text.strip()]
                cell_text = "<br>".join(escape(piece) for piece in pieces) or "&nbsp;"
                classes = (
                    "border border-gray-300 bg-slate-50 px-3 py-2 text-left font-semibold text-gray-800"
                    if cell_tag == "th"
                    else "border border-gray-300 px-3 py-2 align-top text-gray-700"
                )
                cells_html.append(f"<{cell_tag} class='{classes}'>{cell_text}</{cell_tag}>")
            rows_html.append("<tr>" + "".join(cells_html) + "</tr>")

        if rows_html:
            blocks.append(
                "<div class='my-5 overflow-x-auto rounded-xl border border-gray-200 shadow-sm'>"
                "<table class='min-w-full border-collapse bg-white text-sm'>"
                + "".join(rows_html)
                + "</table></div>"
            )

    for child in document.element.body.iterchildren():
        if child.tag == qn("w:p"):
            render_paragraph(Paragraph(child, document))
        elif child.tag == qn("w:tbl"):
            render_table(Table(child, document))

    if not blocks:
        blocks.append("<p class='text-gray-500'>文档中没有可显示的内容。</p>")

    return "html", "".join(blocks)


def is_header_row(row: list[str]) -> bool:
    values = [cell.strip() for cell in row if cell.strip()]
    if not values:
        return False
    if len(values) == 1:
        return False
    return True


def render_xlsx_preview(target: Path) -> tuple[str, str]:
    workbook = load_workbook(filename=str(target), read_only=True, data_only=True)
    sheet_blocks: list[str] = []
    sheet_names = [escape(sheet.title) for sheet in workbook.worksheets[:8]]

    if sheet_names:
        sheet_blocks.append(
            "<div class='mb-5 flex flex-wrap gap-2'>"
            + "".join(
                f"<span class='rounded-full bg-blue-50 px-3 py-1 text-xs font-medium text-blue-700'>{name}</span>"
                for name in sheet_names
            )
            + "</div>"
        )

    for sheet in workbook.worksheets[:8]:
        row_values: list[list[str]] = []
        for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            if row_index > OFFICE_SHEET_ROW_LIMIT:
                break
            values = ["" if cell is None else str(cell) for cell in row[:OFFICE_SHEET_COL_LIMIT]]
            while values and values[-1] == "":
                values.pop()
            if values:
                row_values.append(values)

        max_cols = max((len(row) for row in row_values), default=0)
        has_header = bool(row_values) and is_header_row(row_values[0])
        header = row_values[0] if has_header else [f"列 {index + 1}" for index in range(max_cols)]
        body_rows = row_values[1:] if has_header else row_values

        table_html = []
        table_html.append("<div class='overflow-x-auto rounded-xl border border-gray-200 shadow-sm'>")
        table_html.append("<table class='min-w-full border-collapse bg-white text-sm'>")
        table_html.append("<thead class='bg-slate-50'>")
        table_html.append("<tr>")
        table_html.append("<th class='border border-gray-300 px-3 py-2 text-left text-xs font-semibold uppercase text-slate-500'>#</th>")
        for title in header:
            table_html.append(
                f"<th class='border border-gray-300 px-3 py-2 text-left font-semibold text-gray-800'>{escape(title)}</th>"
            )
        table_html.append("</tr></thead><tbody>")

        for row_index, row in enumerate(body_rows, start=1):
            padded = row + [""] * (len(header) - len(row))
            table_html.append("<tr class='odd:bg-white even:bg-slate-50/60'>")
            table_html.append(
                f"<td class='border border-gray-200 px-3 py-2 text-xs text-slate-400'>{row_index}</td>"
            )
            for cell in padded:
                table_html.append(
                    f"<td class='border border-gray-200 px-3 py-2 align-top text-gray-700'>{escape(cell) or '&nbsp;'}</td>"
                )
            table_html.append("</tr>")

        if not body_rows and not has_header:
            table_html.append("<tr><td class='border border-gray-200 px-3 py-4 text-gray-500' colspan='2'>工作表为空。</td></tr>")

        table_html.append("</tbody></table></div>")

        notes = []
        if sheet.max_row > OFFICE_SHEET_ROW_LIMIT:
            notes.append(f"仅显示前 {OFFICE_SHEET_ROW_LIMIT} 行")
        if sheet.max_column > OFFICE_SHEET_COL_LIMIT:
            notes.append(f"仅显示前 {OFFICE_SHEET_COL_LIMIT} 列")

        sheet_blocks.append(
            "<section class='mb-8'>"
            f"<div class='mb-3 flex items-center justify-between gap-3'>"
            f"<div><h3 class='text-base font-semibold text-gray-900'>{escape(sheet.title)}</h3>"
            f"<p class='text-xs text-gray-400'>共 {sheet.max_row} 行 · {sheet.max_column} 列</p></div>"
            f"<div class='text-xs text-amber-600'>{' · '.join(notes) if notes else ''}</div>"
            "</div>"
            + "".join(table_html)
            + "</section>"
        )

    workbook.close()
    if not sheet_blocks:
        return "html", "<p>工作簿中没有可显示的内容。</p>"
    return "html", "".join(sheet_blocks)


def extract_slide_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs).strip()
                if text:
                    texts.append(text)
        elif hasattr(shape, "text") and shape.text:
            text = shape.text.strip()
            if text:
                texts.append(text)
    return texts


def render_pptx_text_preview(target: Path) -> tuple[str, str]:
    presentation = Presentation(str(target))
    slide_blocks: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts = extract_slide_texts(slide)
        title = texts[0] if texts else f"第 {slide_index} 页"
        bullets = texts[1:] if len(texts) > 1 else []
        preview_body = (
            "<ul class='mt-3 space-y-2 text-sm text-gray-700'>"
            + "".join(
                f"<li class='rounded-lg bg-slate-50 px-3 py-2 leading-6'>{escape(item)}</li>"
                for item in bullets[:8]
            )
            + "</ul>"
            if bullets
            else "<p class='mt-3 text-sm text-gray-500'>此页没有可提取的正文文本。</p>"
        )
        slide_blocks.append(
            "<section class='mb-6 overflow-hidden rounded-2xl border border-gray-200 bg-white shadow-sm'>"
            "<div class='border-b border-gray-100 bg-gradient-to-r from-slate-50 to-white px-5 py-4'>"
            f"<div class='text-xs font-medium uppercase tracking-wide text-blue-600'>Slide {slide_index}</div>"
            f"<h3 class='mt-1 text-lg font-semibold text-gray-900'>{escape(title)}</h3>"
            "</div>"
            f"<div class='px-5 py-4'>{preview_body}</div>"
            "</section>"
        )

    return "html", "".join(slide_blocks or ["<p>演示文稿中没有可显示的内容。</p>"])


def render_ppt_preview(target: Path) -> tuple[str, str]:
    cache_dir = office_cache_dir(target)
    image_dir = cache_dir / "slides"
    image_dir.mkdir(parents=True, exist_ok=True)

    slide_images = sorted(image_dir.glob("slide-*.png"))
    if not slide_images:
        cleanup_cache_dir(image_dir)
        slide_images = export_ppt_slides_with_com(target, image_dir)

    if slide_images:
        image_blocks = []
        cache_key = office_cache_key(target)
        for index, image_path in enumerate(sorted(slide_images), start=1):
            image_url = f"/office-cache/{cache_key}/slides/{image_path.name}"
            image_blocks.append(
                "<section class='mb-8 overflow-hidden rounded-2xl border border-gray-200 bg-white shadow-sm'>"
                "<div class='border-b border-gray-100 bg-gradient-to-r from-slate-50 to-white px-5 py-4'>"
                f"<div class='text-xs font-medium uppercase tracking-wide text-blue-600'>Slide {index}</div>"
                f"<h3 class='mt-1 text-lg font-semibold text-gray-900'>第 {index} 页</h3>"
                "</div>"
                f"<img src='{image_url}' alt='Slide {index}' class='w-full bg-slate-100' />"
                "</section>"
            )
        return "html", "".join(image_blocks)

    return render_pptx_text_preview(target)


def render_office_preview(target: Path, kind: str) -> tuple[str, str, str, str]:
    pdf_url = build_office_pdf_preview(target)
    if pdf_url:
        return "pdf", "", "LibreOffice PDF 预览", pdf_url

    source_path, notice = convert_legacy_office(target, kind)
    if source_path != target:
        pdf_url = build_office_pdf_preview(source_path)
        if pdf_url:
            label = f"{notice} LibreOffice PDF 预览" if notice else "LibreOffice PDF 预览"
            return "pdf", "", label, pdf_url

    if kind == "word":
        mode, html = render_docx_preview(source_path)
    elif kind == "excel":
        mode, html = render_xlsx_preview(source_path)
    elif kind == "powerpoint":
        mode, html = render_ppt_preview(source_path)
    else:
        return "legacy", "<p>暂不支持预览此类 Office 文件。</p>", "", ""

    if notice:
        notice_html = (
            "<div class='mb-4 rounded-xl border border-amber-200 bg-amber-50 px-4 py-3 text-sm text-amber-800'>"
            f"{escape(notice)}</div>"
        )
        html = notice_html + html
    return mode, html, notice or "兼容预览", ""


def build_entry(path: Path, mount_base: Path, mount_name: str) -> dict:
    stat = path.stat()
    is_dir = path.is_dir()
    rel = path.relative_to(mount_base).as_posix()
    ext = path.suffix.lower()
    size = stat.st_size if not is_dir else 0
    return {
        "name": path.name,
        "is_dir": is_dir,
        "size": size,
        "size_str": format_size(size) if not is_dir else "-",
        "modified": stat.st_mtime,
        "ext": ext,
        "kind": "folder" if is_dir else classify(ext),
        "rel_path": rel,
        "browse_url": f"/browse/{mount_name}/{rel}" if is_dir else f"/preview/{mount_name}/{rel}",
        "download_url": f"/download/{mount_name}/{rel}" if not is_dir else None,
    }


@router.get("/", response_class=HTMLResponse)
async def index(request: Request):
    if not is_authenticated(request):
        return RedirectResponse("/login", status_code=302)
    return templates.TemplateResponse(
        request,
        "index.html",
        {
            "mounts": settings.mounts,
            "current_mount": None,
            "entries": [],
            "breadcrumbs": [],
            "current_path": "",
            "mount_name": "",
            "writable": False,
            "csrf_token": get_csrf_token(request),
        },
    )


@router.get("/browse/{mount_name}", response_class=HTMLResponse)
@router.get("/browse/{mount_name}/{dir_path:path}", response_class=HTMLResponse)
async def browse(request: Request, mount_name: str, dir_path: str = ""):
    if not is_authenticated(request):
        return RedirectResponse(f"/login?next=/browse/{mount_name}/{dir_path}", status_code=302)

    mount = get_mount(mount_name)
    if not mount:
        raise HTTPException(404, "挂载点不存在")

    target = safe_resolve(mount.path, dir_path)
    if not target or not target.exists():
        raise HTTPException(404, "路径不存在")

    if target.is_file():
        return RedirectResponse(f"/preview/{mount_name}/{dir_path}", status_code=302)

    try:
        raw_entries = list(target.iterdir())
    except PermissionError as exc:
        raise HTTPException(403, "没有权限访问此目录") from exc

    mount_base = Path(mount.path).resolve()
    entries = []
    for item in raw_entries:
        try:
            entries.append(build_entry(item, mount_base, mount_name))
        except (OSError, ValueError):
            continue

    entries.sort(key=lambda entry: (not entry["is_dir"], entry["name"].lower()))

    breadcrumbs = [{"name": mount_name, "url": f"/browse/{mount_name}"}]
    parts = [part for part in dir_path.split("/") if part]
    for index, part in enumerate(parts):
        crumb_path = "/".join(parts[: index + 1])
        breadcrumbs.append({"name": part, "url": f"/browse/{mount_name}/{crumb_path}"})

    return templates.TemplateResponse(
        request,
        "index.html",
        {
            "mounts": settings.mounts,
            "current_mount": mount,
            "entries": entries,
            "breadcrumbs": breadcrumbs,
            "current_path": dir_path,
            "mount_name": mount_name,
            "writable": mount.writable,
            "upload_url": f"/upload/{mount_name}/{dir_path}" if dir_path else f"/upload/{mount_name}",
            "csrf_token": get_csrf_token(request),
        },
    )


@router.get("/preview/{mount_name}/{file_path:path}", response_class=HTMLResponse)
async def preview(request: Request, mount_name: str, file_path: str):
    if not is_authenticated(request):
        return RedirectResponse(f"/login?next=/preview/{mount_name}/{file_path}", status_code=302)

    mount = get_mount(mount_name)
    if not mount:
        raise HTTPException(404)

    target = safe_resolve(mount.path, file_path)
    if not target or not target.is_file():
        raise HTTPException(404, "文件不存在")

    ext = target.suffix.lower()
    kind = classify(ext)

    parts = [part for part in file_path.split("/") if part]
    parent_path = "/".join(parts[:-1])
    parent_url = f"/browse/{mount_name}/{parent_path}" if parent_path else f"/browse/{mount_name}"

    text_content = ""
    office_preview_mode = ""
    office_preview_html = ""
    office_preview_label = ""
    office_preview_url = ""
    if kind == "text":
        try:
            size = target.stat().st_size
            if size <= TEXT_PREVIEW_LIMIT:
                text_content = target.read_text(encoding="utf-8", errors="replace")
            else:
                text_content = f"[文件过大（{format_size(size)}），无法预览，请下载后查看]"
        except Exception as exc:  # noqa: BLE001
            text_content = f"[读取失败: {exc}]"
    elif kind in {"word", "excel", "powerpoint"}:
        try:
            office_preview_mode, office_preview_html, office_preview_label, office_preview_url = render_office_preview(target, kind)
        except Exception as exc:  # noqa: BLE001
            office_preview_mode = "error"
            office_preview_html = f"<p>预览生成失败：{escape(str(exc))}</p>"

    return templates.TemplateResponse(
        request,
        "preview.html",
        {
            "mounts": settings.mounts,
            "mount_name": mount_name,
            "file_path": file_path,
            "file_name": target.name,
            "file_size": format_size(target.stat().st_size),
            "kind": kind,
            "ext": ext,
            "text_content": text_content,
            "office_preview_mode": office_preview_mode,
            "office_preview_html": office_preview_html,
            "office_preview_label": office_preview_label,
            "office_preview_url": office_preview_url,
            "parent_url": parent_url,
            "download_url": f"/download/{mount_name}/{file_path}",
            "raw_url": f"/raw/{mount_name}/{file_path}",
            "csrf_token": get_csrf_token(request),
        },
    )


@router.get("/office-cache/{cache_key}/{sub_path:path}")
async def office_cache_file(request: Request, cache_key: str, sub_path: str):
    if not is_authenticated(request):
        raise HTTPException(401)

    target = safe_resolve(str(OFFICE_CACHE_DIR / cache_key), sub_path)
    if not target or not target.is_file():
        raise HTTPException(404)

    mime, _ = mimetypes.guess_type(str(target))
    return FileResponse(path=target, media_type=mime or "application/octet-stream")


@router.get("/download/{mount_name}/{file_path:path}")
async def download(request: Request, mount_name: str, file_path: str):
    if not is_authenticated(request):
        raise HTTPException(401, "请先登录")

    mount = get_mount(mount_name)
    if not mount:
        raise HTTPException(404)

    target = safe_resolve(mount.path, file_path)
    if not target or not target.is_file():
        raise HTTPException(404, "文件不存在")

    return FileResponse(path=target, filename=target.name, media_type="application/octet-stream")


@router.get("/raw/{mount_name}/{file_path:path}")
async def raw(request: Request, mount_name: str, file_path: str):
    if not is_authenticated(request):
        raise HTTPException(401)

    mount = get_mount(mount_name)
    if not mount:
        raise HTTPException(404)

    target = safe_resolve(mount.path, file_path)
    if not target or not target.is_file():
        raise HTTPException(404)

    mime, _ = mimetypes.guess_type(str(target))
    return FileResponse(path=target, media_type=mime or "application/octet-stream")


@router.post("/upload/{mount_name}")
@router.post("/upload/{mount_name}/{dir_path:path}")
async def upload(request: Request, mount_name: str, dir_path: str = "", file: UploadFile = File(...)):
    if not is_authenticated(request):
        raise HTTPException(401, "请先登录")
    validate_csrf(request)

    mount = get_mount(mount_name)
    if not mount:
        raise HTTPException(404, "挂载点不存在")
    if not mount.writable:
        raise HTTPException(403, "该挂载点为只读")

    target_dir = safe_resolve(mount.path, dir_path)
    if not target_dir or not target_dir.is_dir():
        raise HTTPException(404, "目标目录不存在")

    filename = Path(file.filename).name if file.filename else ""
    if not filename:
        raise HTTPException(400, "文件名无效")

    dest = target_dir / filename
    content = await file.read()
    dest.write_bytes(content)

    return JSONResponse({"success": True, "filename": filename, "size": len(content)})


@router.delete("/delete/{mount_name}/{file_path:path}")
async def delete_file(request: Request, mount_name: str, file_path: str):
    if not is_authenticated(request):
        raise HTTPException(401)
    validate_csrf(request)

    mount = get_mount(mount_name)
    if not mount:
        raise HTTPException(404)
    if not mount.writable:
        raise HTTPException(403, "该挂载点为只读")

    target = safe_resolve(mount.path, file_path)
    if not target or not target.exists():
        raise HTTPException(404)
    if target == Path(mount.path).resolve():
        raise HTTPException(400, "不能删除挂载根目录")

    if target.is_file():
        target.unlink()
    elif target.is_dir():
        shutil.rmtree(target)

    return JSONResponse({"success": True})
